import streamlit as st
import google.generativeai as genai
import pandas as pd
import docx
import PyPDF2
import io
import json
import openpyxl
import markdown
import re
from bs4 import BeautifulSoup
import ebooklib
from ebooklib import epub
import pptx
from io import BytesIO

# Set page configuration
st.set_page_config(
    page_title="Document Chat AI",
    page_icon="📑",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for modern UI
st.markdown("""
<style>
    .main {
        background-color: #032c54;
    }
    .stApp {
        max-width: 1200px;
        margin: 0 auto;
    }
    .chat-message {
        padding: 1.5rem;
        border-radius: 0.5rem;
        margin-bottom: 1rem;
        display: flex;
        align-items: flex-start;
    }
    .chat-message.user {
        background-color: #032c54;
    }
    .chat-message.bot {
        background-color: #dce8f7;
    }
    .chat-message .avatar {
        width: 40px;
        height: 40px;
        border-radius: 50%;
        object-fit: cover;
        margin-right: 1rem;
    }
    .chat-message .message {
        flex-grow: 1;
    }
    .file-list {
        background-color: #f1f3f5;
        padding: 1rem;
        border-radius: 0.5rem;
        margin-bottom: 1rem;
    }
    .file-item {
        display: flex;
        align-items: center;
        padding: 0.5rem;
        border-bottom: 1px solid #dee2e6;
    }
    .file-icon {
        margin-right: 0.5rem;
        font-size: 1.2rem;
    }
    .stButton>button {
        background-color: #4e73df;
        color: white;
        border-radius: 0.3rem;
        border: none;
        padding: 0.5rem 1rem;
        font-weight: 500;
    }
    .stButton>button:hover {
        background-color: #375bd2;
    }
    .sidebar-header {
        padding: 1rem;
        background-color: #4e73df;
        color: white;
        border-radius: 0.5rem;
        margin-bottom: 1rem;
    }
    .feedback-buttons button {
        background-color: transparent;
        border: 1px solid #dee2e6;
        border-radius: 0.3rem;
        padding: 0.3rem 0.7rem;
        margin-right: 0.5rem;
    }
    .stTextInput>div>div>input {
        border-radius: 0.3rem;
    }
    h1, h2, h3 {
        color: #3a3b45;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state variables
if "api_key" not in st.session_state:
    st.session_state.api_key = None
if "model" not in st.session_state:
    st.session_state.model = "gemini-2.0-flash"
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "files_content" not in st.session_state:
    st.session_state.files_content = {}
if "last_uploaded_files" not in st.session_state:
    st.session_state.last_uploaded_files = []
if "current_tab" not in st.session_state:
    st.session_state.current_tab = "Chat"

# File icons based on type
FILE_ICONS = {
    "PDF": "📄",
    "DOCX": "📝",
    "TXT": "📃",
    "CSV": "📊",
    "XLSX": "📊",
    "JSON": "🔢",
    "MD": "📑",
    "PPTX": "📊",
    "HTML": "🌐",
    "EPUB": "📚"
}

# File extraction functions
def extract_text_from_pdf(file):
    pdf_reader = PyPDF2.PdfReader(file)
    text = ""
    for page_num, page in enumerate(pdf_reader.pages):
        text += f"--- Page {page_num + 1} ---\n"
        text += page.extract_text() + "\n\n"
    return text

def extract_text_from_docx(file):
    doc = docx.Document(file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_txt(file):
    return file.getvalue().decode("utf-8")

def extract_text_from_csv(file):
    df = pd.read_csv(file)
    return df.to_string()

def extract_text_from_xlsx(file):
    excel_data = BytesIO(file.getvalue())
    workbook = openpyxl.load_workbook(excel_data, data_only=True)
    text = ""
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        text += f"Sheet: {sheet_name}\n"
        for row in sheet.iter_rows(values_only=True):
            text += str(row) + "\n"
        text += "\n"
    return text

def extract_text_from_json(file):
    try:
        json_data = json.loads(file.getvalue().decode("utf-8"))
        return json.dumps(json_data, indent=2)
    except:
        return "Error parsing JSON file"

def extract_text_from_markdown(file):
    md_content = file.getvalue().decode("utf-8")
    # Convert to HTML and then extract text
    html = markdown.markdown(md_content)
    soup = BeautifulSoup(html, 'html.parser')
    return soup.get_text()

def extract_text_from_pptx(file):
    prs = pptx.Presentation(file)
    text = ""
    for i, slide in enumerate(prs.slides):
        text += f"--- Slide {i+1} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        text += "\n"
    return text

def extract_text_from_html(file):
    html_content = file.getvalue().decode("utf-8")
    soup = BeautifulSoup(html_content, 'html.parser')
    return soup.get_text()

def extract_text_from_epub(file):
    book = epub.read_epub(file)
    text = ""
    for item in book.get_items():
        if item.get_type() == ebooklib.ITEM_DOCUMENT:
            soup = BeautifulSoup(item.get_content(), 'html.parser')
            text += soup.get_text() + "\n\n"
    return text

# Process file based on type
def process_file(uploaded_file):
    try:
        file_type = uploaded_file.type
        file_extension = uploaded_file.name.split('.')[-1].lower()
        
        if file_type == "application/pdf" or file_extension == "pdf":
            return extract_text_from_pdf(uploaded_file), "PDF"
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or file_extension == "docx":
            return extract_text_from_docx(uploaded_file), "DOCX"
        elif file_type == "text/plain" or file_extension == "txt":
            return extract_text_from_txt(uploaded_file), "TXT"
        elif file_type == "text/csv" or file_extension == "csv":
            return extract_text_from_csv(uploaded_file), "CSV"
        elif file_type == "application/json" or file_extension == "json":
            return extract_text_from_json(uploaded_file), "JSON"
        elif "spreadsheetml" in file_type or file_extension == "xlsx":
            return extract_text_from_xlsx(uploaded_file), "XLSX"
        elif file_extension == "md" or file_extension == "markdown":
            return extract_text_from_markdown(uploaded_file), "MD"
        elif "presentationml" in file_type or file_extension == "pptx":
            return extract_text_from_pptx(uploaded_file), "PPTX"
        elif file_type == "text/html" or file_extension == "html" or file_extension == "htm":
            return extract_text_from_html(uploaded_file), "HTML"
        elif file_extension == "epub":
            return extract_text_from_epub(uploaded_file), "EPUB"
        else:
            return None, "Unsupported"
    except Exception as e:
        return f"Error processing file: {str(e)}", "Error"

# Generate AI response
def generate_ai_response(prompt, files_content, model_name, api_key):
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel(model_name)
        
        # Combine file contents with truncation to avoid token limits
        combined_content = ""
        for filename, content in files_content.items():
            file_content = content[:7500]  # Limit each file content
            combined_content += f"[File: {filename}]\n{file_content}\n\n"
        
        # Generate response with clear instructions
        response = model.generate_content(
            f"""Files content:
{combined_content}

User question: {prompt}

IMPORTANT INSTRUCTIONS:
1. Base your response ONLY on the files content provided above
2. Reference specific file names when providing information
3. If the information is not in the files, clearly state that
4. Format your response with markdown for better readability
5. DO NOT reference any previous conversations or files not listed above
6. Be concise but thorough in your answers"""
        )
        
        return response.text
    except Exception as e:
        return f"Error generating response: {str(e)}"

# Sidebar configuration
with st.sidebar:
    st.markdown("""
    <div class="sidebar-header">
        <h1 style="margin:0; font-size:1.5rem;">📑 Document Chat AI</h1>
    </div>
    """, unsafe_allow_html=True)
    
    # Tabs
    tab_options = ["Chat", "Settings", "About"]
    tabs = st.radio("", tab_options, horizontal=True, 
                   index=tab_options.index(st.session_state.current_tab))
    st.session_state.current_tab = tabs
    
    # API Configuration
    if st.session_state.current_tab == "Settings":
        st.markdown("### 🔑 API Configuration")
        api_key = st.text_input("Enter Google Gemini API Key:", type="password", 
                             help="Get your key from [Google AI Studio](https://aistudio.google.com/app/apikey)")
        if api_key:
            st.session_state.api_key = api_key
        
        # Model Settings
        st.markdown("### ⚙️ Model Settings")
        model_options = [
            "gemini-2.0-flash", 
            "gemini-1.5-flash", 
            "gemini-1.5-pro", 
            "gemini-2.0-pro-exp-02-05",
            "gemini-2.5-pro-exp-03-25", 
            "gemini-1.5-flash-8b"
        ]
        st.session_state.model = st.selectbox("Select Gemini Model:", model_options, 
                                           index=model_options.index(st.session_state.model))
        
        # Advanced Settings
        st.markdown("### 🔧 Advanced Settings")
        st.slider("Response Temperature", min_value=0.0, max_value=1.0, value=0.2, step=0.1, 
                 help="Lower values make output more focused and deterministic")
        st.slider("Maximum Tokens", min_value=100, max_value=8000, value=4000, step=100, 
                 help="Maximum length of the generated response")
    
    # About section
    elif st.session_state.current_tab == "About":
        st.markdown("### 📚 About Document Chat AI")
        st.markdown("""
        Document Chat AI allows you to chat with your files using Google's Gemini API.
        
        **Supported file types:**
        - PDF (.pdf) - Documents, reports, papers
        - Word (.docx) - Text documents
        - Text (.txt) - Plain text files
        - CSV (.csv) - Tabular data
        - Excel (.xlsx) - Spreadsheets
        - JSON (.json) - Structured data
        - Markdown (.md) - Formatted text
        - PowerPoint (.pptx) - Presentations
        - HTML (.html) - Web pages
        - EPUB (.epub) - E-books
        
        **Privacy notice:**
        Files are processed locally in your browser and not stored permanently on any server.
        
        **Need help?**
        Contact support at [example@email.com](mailto:example@email.com)
        """)
    
    # Files section always visible
    if st.session_state.current_tab == "Chat":
        st.markdown("### 📂 Uploaded Documents")
        if st.session_state.files_content:
            for filename, content in st.session_state.files_content.items():
                file_ext = filename.split('.')[-1].upper()
                icon = FILE_ICONS.get(file_ext, "📄")
                st.markdown(f"""
                <div class="file-item">
                    <span class="file-icon">{icon}</span>
                    <span>{filename}</span>
                </div>
                """, unsafe_allow_html=True)
        else:
            st.info("No documents uploaded yet")
        
        if st.button("Clear All Files", key="clear_sidebar", use_container_width=True):
            st.session_state.files_content = {}
            st.session_state.chat_history = []
            st.session_state.last_uploaded_files = []
            st.rerun()

# Main content
if st.session_state.current_tab == "Chat":
    st.markdown("## 📑 Document Chat AI")
    st.markdown("Upload your documents and ask questions about their content")
    
    # File upload section
    col1, col2 = st.columns([3, 1])
    
    with col1:
        uploaded_files = st.file_uploader(
            "Upload your documents",
            type=["pdf", "docx", "txt", "csv", "xlsx", "json", "md", "pptx", "html", "epub"],
            accept_multiple_files=True
        )
    with col2:
        st.markdown("&nbsp;")  # Add space for alignment
        if st.button("Process Files", use_container_width=True):
            if uploaded_files:
                # Files changed, clear previous files
                current_files = [file.name for file in uploaded_files]
                if current_files != st.session_state.last_uploaded_files:
                    st.session_state.files_content = {}
                    st.session_state.last_uploaded_files = current_files
                
                with st.spinner("Processing your documents..."):
                    # Process each uploaded file
                    for uploaded_file in uploaded_files:
                        # Check if file was already processed
                        if uploaded_file.name not in st.session_state.files_content:
                            extracted_text, file_type = process_file(uploaded_file)
                            if extracted_text and file_type != "Error" and file_type != "Unsupported":
                                st.session_state.files_content[uploaded_file.name] = extracted_text
                                icon = FILE_ICONS.get(file_type, "📄")
                                st.success(f"{icon} {file_type} file '{uploaded_file.name}' processed successfully!")
                            else:
                                st.error(f"❌ Could not process '{uploaded_file.name}': {extracted_text}")
    
    # Display document summary if files are uploaded
    if st.session_state.files_content:
        file_count = len(st.session_state.files_content)
        st.markdown(f"""
        <div class="file-list">
            <h3 style="color: #1E90FF;">📚 {file_count} Document{'' if file_count == 1 else 's'} Processed</h3>
             <p style="color: #1E90FF;"You can now ask questions about your documents.</p>
        </div>
        """, unsafe_allow_html=True)
    
    # Chat interface
    st.markdown("### 💬 Chat with your documents")
    
    # Display chat history
    chat_container = st.container()
    with chat_container:
        for message in st.session_state.chat_history:
            if message["role"] == "user":
                st.markdown(f"""
                <div class="chat-message user">
                    <img src="https://avatars.githubusercontent.com/u/0" class="avatar">
                    <div class="message">
                        <b>You:</b><br>{message['content']}
                    </div>
                </div>
                """, unsafe_allow_html=True)
            else:
                st.markdown(f"""
                <div class="chat-message bot">
                    <img src="https://avatars.githubusercontent.com/u/1" class="avatar">
                    <div class="message">
                        <b>AI Assistant:</b><br>{message['content']}
                    </div>
                </div>
                """, unsafe_allow_html=True)

    
    # User input
    st.markdown("#### Ask a question")
    col1, col2 = st.columns([4, 1])
    
    with col1:
        user_question = st.text_input("", placeholder="Type your question here...", key="user_input", label_visibility="collapsed")
    with col2:
        submit_button = st.button("Send 📩", use_container_width=True)
    
    # Action buttons for chat interface
    col1, col2 = st.columns([1, 1])
    with col1:
        if st.button("Clear Chat", use_container_width=True, key="clear_chat"):
            st.session_state.chat_history = []
            st.rerun()
    with col2:
        if st.button("Sample Questions", use_container_width=True, key="sample_questions"):
            st.info("""
            Try asking:
            - What is the main topic of these documents?
            - Summarize the content of the PDF file
            - What data is in the CSV file?
            - Compare the information between documents
            """)
    
    # Process user question
    if submit_button and user_question:
        if not st.session_state.api_key:
            st.error("Please enter your Gemini API key in the Settings tab.")
        elif not st.session_state.files_content:
            st.warning("Please upload and process at least one document first.")
        else:
            # Add user message to chat history
            st.session_state.chat_history.append({"role": "user", "content": user_question})
            
            # Generate AI response
            with st.spinner("AI is analyzing your documents..."):
                ai_response = generate_ai_response(
                    user_question, 
                    st.session_state.files_content, 
                    st.session_state.model, 
                    st.session_state.api_key
                )
                
                # Add AI response to chat history
                st.session_state.chat_history.append({"role": "assistant", "content": ai_response})
            
            # Rerun to update chat display
            st.rerun()
    
    # Display message if no files uploaded
    if not st.session_state.files_content:
        st.info("👆 Please upload your documents to start chatting with them")
    
    # Tips for better results
    with st.expander("💡 Tips for better results"):
        st.markdown("""
        - Ask specific questions about the content in your documents
        - When working with multiple files, specify which file you're asking about
        - For large documents, try to reference specific sections or topics
        - For data files (CSV, Excel), you can ask for specific data analysis
        - Try asking for summaries, comparisons, or key insights from your documents
        """)

# Footer
st.markdown("---")
st.markdown("""
<p style="text-align: center; color: #6c757d; font-size: 0.8rem;">
    Document Chat AI | Built with Streamlit and Google Gemini | © 2025
</p>
""", unsafe_allow_html=True)
